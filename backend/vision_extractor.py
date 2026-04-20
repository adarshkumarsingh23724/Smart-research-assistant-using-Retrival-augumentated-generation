"""
vision_extractor.py
───────────────────
Multimodal image extraction and AI captioning for the Smart Research Assistant.

Responsibilities:
  1. Extract images/diagrams from PDF pages (via PyMuPDF / fitz)
  2. Extract images/diagrams from PPTX slides (via python-pptx)
  3. Send each image to Groq's vision LLM and produce a rich academic caption
     (parallel processing via ThreadPoolExecutor for speed)
  4. Return LangChain Document objects ready to be added to ChromaDB

The `ENABLE_VISION` env var acts as a kill-switch:
    ENABLE_VISION=false  →  all functions return [] immediately (no-op, safe)
"""

import os
import io
import base64
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Optional, Tuple

from dotenv import load_dotenv
from langchain_core.documents import Document

load_dotenv()

# ──────────────────────────────────────────────
# Configuration
# ──────────────────────────────────────────────
ENABLE_VISION: bool = os.getenv("ENABLE_VISION", "true").lower() == "true"
VISION_MODEL: str = os.getenv(
    "VISION_MODEL",
    "meta-llama/llama-4-scout-17b-16e-instruct"
)
MIN_IMAGE_DIMENSION: int = 120      # skip images < this (px) — filters icons/bullets
MAX_IMAGES_PER_DOC:  int = 20       # free-tier safe cap (≈ 10-11 imgs/min at 30k TPM)
VISION_WORKERS:      int = 2        # 2 parallel workers — avoids TPM burst on free tier
VISION_MAX_RETRIES:  int = 3        # how many times to retry a 429 before giving up

# Runtime kill-switch — set on fatal errors (bad model/key). Resets per process.
_vision_runtime_disabled: bool = False

# Academic-focused vision prompt (concise for speed, still thorough)
VISION_PROMPT = """You are an expert academic assistant analyzing content from a lecture slide or educational document.

Examine this image and provide a concise but complete academic description:

1. **Image Type**: diagram, flowchart, graph, table, code snippet, equation, photo, or architecture diagram
2. **Content**: Describe all visible details. For graphs: axes, labels, trends, key values. For diagrams: components, arrows, labels, relationships. For tables: columns, rows, key data. For code: transcribe exactly. For equations: write clearly.
3. **Key Insight**: The main academic concept or takeaway from this image.

Be precise — your description will be used to answer exam questions about this visual content."""


# Sentinel exception for rate-limit errors — allows the retry wrapper to parse wait time
class _RateLimitError(Exception):
    pass


def _is_vision_available() -> bool:
    """Check if vision processing is enabled and the API key is present."""
    global _vision_runtime_disabled
    if _vision_runtime_disabled:
        return False
    if not ENABLE_VISION:
        return False
    api_key = os.getenv("GROQ_API_KEY", "").strip()
    return bool(api_key)


def _pil_image_to_base64(pil_image) -> str:
    """Convert a PIL Image to a base64-encoded JPEG string."""
    buffer = io.BytesIO()
    if pil_image.mode not in ("RGB", "L"):
        pil_image = pil_image.convert("RGB")
    pil_image.save(buffer, format="JPEG", quality=80)
    buffer.seek(0)
    return base64.b64encode(buffer.read()).decode("utf-8")


def _describe_image_with_vision(pil_image, context_hint: str = "") -> Optional[str]:
    """
    Send a PIL Image to the Groq vision LLM and return a rich academic caption.
    Returns None if the call fails or vision is disabled.
    On fatal errors (decommissioned model, bad key), sets _vision_runtime_disabled=True
    so no further calls are wasted in this ingest run.
    """
    global _vision_runtime_disabled

    if not _is_vision_available():
        return None

    try:
        from langchain_groq import ChatGroq
        from langchain_core.messages import HumanMessage

        load_dotenv(override=True)
        api_key = os.getenv("GROQ_API_KEY", "").strip()

        llm = ChatGroq(
            api_key=api_key,
            model=VISION_MODEL,
            temperature=0,
            max_tokens=512,   # compact output — fits ~11 images in free-tier 30k TPM/min
        )

        image_b64 = _pil_image_to_base64(pil_image)
        context_text = f"\n\nAdditional context: {context_hint}" if context_hint else ""

        message = HumanMessage(
            content=[
                {"type": "text",      "text": VISION_PROMPT + context_text},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_b64}"}}
            ]
        )

        response = llm.invoke([message])
        return response.content.strip() if response.content else None

    except Exception as e:
        err_str = str(e)
        err_lower = err_str.lower()
        if "model_decommissioned" in err_lower or "decommissioned" in err_lower:
            _vision_runtime_disabled = True
            print(f"[Vision] FATAL: Vision model '{VISION_MODEL}' has been decommissioned.")
            print(f"[Vision] Update VISION_MODEL in .env — recommended: meta-llama/llama-4-scout-17b-16e-instruct")
            print(f"[Vision] Disabling vision for the remainder of this ingest run.")
        elif "invalid_api_key" in err_lower or "authentication" in err_lower:
            _vision_runtime_disabled = True
            print(f"[Vision] FATAL: Invalid Groq API key. Disabling vision for this run.")
        elif "rate_limit_exceeded" in err_lower or "429" in err_str or "rate limit" in err_lower:
            # Raise so the retry wrapper can catch it and sleep the right amount
            raise _RateLimitError(err_str)
        else:
            print(f"[Vision] API call failed: {e}")
        return None


def _is_image_large_enough(pil_image) -> bool:
    """Return True if the image passes the minimum dimension threshold."""
    w, h = pil_image.size
    return w >= MIN_IMAGE_DIMENSION and h >= MIN_IMAGE_DIMENSION


import re as _re

def _parse_retry_after(err_str: str) -> float:
    """
    Extract the 'retry after N seconds' hint from a Groq 429 error message.
    Returns the wait time in seconds, defaulting to 5.0 if not found.
    """
    # Groq 429 messages contain phrases like:
    #  'Please try again in 1.934s.'
    #  'Please try again in 86ms.'
    #  'Please try again in 4.846s.'
    match_s  = _re.search(r'try again in (\d+(?:\.\d+)?)s', err_str)
    match_ms = _re.search(r'try again in (\d+(?:\.\d+)?)ms', err_str)
    if match_s:
        return float(match_s.group(1)) + 0.5   # add a small safety margin
    if match_ms:
        return float(match_ms.group(1)) / 1000 + 0.5
    return 5.0  # safe default


def _describe_image_with_retry(pil_image, context_hint: str = "") -> Optional[str]:
    """
    Wraps _describe_image_with_vision with smart 429 retry:
    - Parses the exact 'retry after Xs' from the Groq error message
    - Sleeps that long + a 0.5s safety margin, then retries
    - Gives up after VISION_MAX_RETRIES total attempts
    """
    for attempt in range(1, VISION_MAX_RETRIES + 1):
        try:
            result = _describe_image_with_vision(pil_image, context_hint)
            return result   # success (or fatal error that set _vision_runtime_disabled)
        except _RateLimitError as rle:
            if _vision_runtime_disabled:
                return None
            wait = _parse_retry_after(str(rle))
            if attempt < VISION_MAX_RETRIES:
                print(f"[Vision] Rate limit hit. Waiting {wait:.1f}s then retrying (attempt {attempt}/{VISION_MAX_RETRIES})...")
                time.sleep(wait)
            else:
                print(f"[Vision] Rate limit hit on all {VISION_MAX_RETRIES} attempts. Skipping this image.")
    return None


def _caption_task(args: Tuple) -> Tuple[int, int, Optional[str]]:
    """Worker function: caption one image with retry on 429. Returns (location, img_idx, caption)."""
    location, img_idx, pil_img, context_hint = args
    caption = _describe_image_with_retry(pil_img, context_hint)
    return location, img_idx, caption


def extract_images_from_pdf_page(pdf_page) -> List:
    """Extract all qualifying images from a single PyMuPDF page."""
    from PIL import Image
    pil_images = []
    try:
        image_list = pdf_page.get_images(full=True)
        doc = pdf_page.parent
        for img_info in image_list:
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                image_bytes = base_image.get("image", b"")
                if not image_bytes:
                    continue
                pil_img = Image.open(io.BytesIO(image_bytes))
                if _is_image_large_enough(pil_img):
                    pil_images.append(pil_img)
            except Exception:
                pass
    except Exception as e:
        print(f"[Vision] Error scanning PDF page for images: {e}")
    return pil_images


def process_pdf_images(pdf_path: str, source_name: str, course_id: str = None, user_id: str = None) -> List[Document]:
    """
    Extract and AI-caption all images from a PDF using parallel workers.
    Capped at MAX_IMAGES_PER_DOC images total for speed.
    """
    if not _is_vision_available():
        print("[Vision] Skipping PDF vision extraction — disabled or no API key.")
        return []

    try:
        import fitz
    except ImportError:
        print("[Vision] PyMuPDF (fitz) not installed. Run: pip install pymupdf")
        return []

    # ── Phase 1: gather all qualifying images across all pages ────────────────
    all_tasks = []   # list of (page_num, img_idx, pil_img, context_hint)
    try:
        pdf_doc = fitz.open(pdf_path)
        total_pages = len(pdf_doc)
        print(f"[Vision] Scanning {total_pages} PDF page(s) in '{source_name}' for images...")

        for page_num in range(total_pages):
            if len(all_tasks) >= MAX_IMAGES_PER_DOC:
                print(f"[Vision] Hit cap of {MAX_IMAGES_PER_DOC} images. Remaining pages skipped.")
                break
            page = pdf_doc[page_num]
            pil_images = extract_images_from_pdf_page(page)
            for img_idx, pil_img in enumerate(pil_images, start=1):
                if len(all_tasks) >= MAX_IMAGES_PER_DOC:
                    break
                context_hint = f"From page {page_num + 1} of '{source_name}'"
                all_tasks.append((page_num + 1, img_idx, pil_img, context_hint))

        pdf_doc.close()
    except Exception as e:
        print(f"[Vision] Error opening PDF '{source_name}': {e}")
        return []

    if not all_tasks:
        return []

    print(f"[Vision] Captioning {len(all_tasks)} image(s) from '{source_name}' using {VISION_WORKERS} parallel workers...")

    # ── Phase 2: parallel captioning ──────────────────────────────────────────
    docs = []
    worker_args = [(loc, idx, img, hint) for loc, idx, img, hint in all_tasks]

    with ThreadPoolExecutor(max_workers=VISION_WORKERS) as executor:
        futures = {executor.submit(_caption_task, args): args for args in worker_args}
        for future in as_completed(futures):
            if _vision_runtime_disabled:
                executor.shutdown(wait=False, cancel_futures=True)
                break
            try:
                page_num, img_idx, caption = future.result()
                if caption:
                    docs.append(Document(
                        page_content=f"[IMAGE DESCRIPTION - Page {page_num}, Image {img_idx}]\n{caption}",
                        metadata={
                            "source": source_name,
                            "page": page_num,
                            "type": "image_caption",
                            "image_index": img_idx,
                            "course_id": course_id or "unknown",
                            "user_id": user_id or "unknown",
                        }
                    ))
                    print(f"[Vision] Captioned page {page_num}, image {img_idx}.")
            except Exception as e:
                print(f"[Vision] Worker error: {e}")

    print(f"[Vision] PDF scan complete: {len(docs)} caption(s) generated for '{source_name}'.")
    return docs


# ──────────────────────────────────────────────
# PPTX Image Extraction (via python-pptx)
# ──────────────────────────────────────────────

def extract_images_from_pptx_slide(slide) -> List:
    """Extract all qualifying picture shapes from a python-pptx slide."""
    from PIL import Image
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    pil_images = []
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pil_img = Image.open(io.BytesIO(shape.image.blob))
                if _is_image_large_enough(pil_img):
                    pil_images.append(pil_img)
        except Exception:
            pass
    return pil_images


def process_pptx_images(pptx_path: str, source_name: str, course_id: str = None, user_id: str = None) -> List[Document]:
    """
    Extract and AI-caption all images from a PPTX using parallel workers.
    Capped at MAX_IMAGES_PER_DOC images total for speed.
    """
    if not _is_vision_available():
        print("[Vision] Skipping PPTX vision extraction — disabled or no API key.")
        return []

    try:
        from pptx import Presentation
    except ImportError:
        print("[Vision] python-pptx not installed.")
        return []

    # ── Phase 1: gather all qualifying images across all slides ───────────────
    all_tasks = []
    try:
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
        print(f"[Vision] Scanning {total_slides} slide(s) in '{source_name}' for images...")

        for slide_num, slide in enumerate(prs.slides, start=1):
            if len(all_tasks) >= MAX_IMAGES_PER_DOC:
                print(f"[Vision] Hit cap of {MAX_IMAGES_PER_DOC} images. Remaining slides skipped.")
                break
            pil_images = extract_images_from_pptx_slide(slide)
            for img_idx, pil_img in enumerate(pil_images, start=1):
                if len(all_tasks) >= MAX_IMAGES_PER_DOC:
                    break
                context_hint = f"From slide {slide_num} of '{source_name}'"
                all_tasks.append((slide_num, img_idx, pil_img, context_hint))
    except Exception as e:
        print(f"[Vision] Error opening PPTX '{source_name}': {e}")
        return []

    if not all_tasks:
        return []

    print(f"[Vision] Captioning {len(all_tasks)} image(s) from '{source_name}' using {VISION_WORKERS} parallel workers...")

    # ── Phase 2: parallel captioning ──────────────────────────────────────────
    docs = []
    worker_args = [(loc, idx, img, hint) for loc, idx, img, hint in all_tasks]

    with ThreadPoolExecutor(max_workers=VISION_WORKERS) as executor:
        futures = {executor.submit(_caption_task, args): args for args in worker_args}
        for future in as_completed(futures):
            if _vision_runtime_disabled:
                executor.shutdown(wait=False, cancel_futures=True)
                break
            try:
                slide_num, img_idx, caption = future.result()
                if caption:
                    docs.append(Document(
                        page_content=f"[IMAGE DESCRIPTION - Slide {slide_num}, Image {img_idx}]\n{caption}",
                        metadata={
                            "source": source_name,
                            "page": slide_num,
                            "type": "image_caption",
                            "image_index": img_idx,
                            "course_id": course_id or "unknown",
                            "user_id": user_id or "unknown",
                        }
                    ))
                    print(f"[Vision] Captioned slide {slide_num}, image {img_idx}.")
            except Exception as e:
                print(f"[Vision] Worker error: {e}")

    print(f"[Vision] PPTX scan complete: {len(docs)} caption(s) generated for '{source_name}'.")
    return docs


