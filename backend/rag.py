import os
import tempfile
from tempfile import NamedTemporaryFile
from langchain_groq import ChatGroq
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_core.documents import Document

from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from dotenv import load_dotenv

load_dotenv()

# Initialize components
embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

CHROMA_DB_PATH = "./chroma_db"

def get_vector_store():
    global vector_store
            
    # Ensure directory exists before initializing Chroma
    os.makedirs(CHROMA_DB_PATH, exist_ok=True)
    
    # Initialize or load Chroma DB
    try:
        vector_store = Chroma(
            persist_directory=CHROMA_DB_PATH,
            embedding_function=embeddings
        )
    except Exception as e:
        print(f"Error loading Chroma DB: {e}")
        # Make a new one
        vector_store = Chroma(
            embedding_function=embeddings,
            persist_directory=CHROMA_DB_PATH
        )
        
    return vector_store

vector_store = None

# Supported Groq models — sourced from https://console.groq.com/docs/models
# Only active (non-deprecated) IDs are listed here.
GROQ_MODELS = {
    # Production models (stable, recommended for production use)
    "llama-3.3-70b-versatile":               "Llama 3.3 70B (default)",
    "llama-3.1-8b-instant":                  "Llama 3.1 8B (ultra fast)",
    "openai/gpt-oss-120b":                   "GPT-OSS 120B (OpenAI)",
    "openai/gpt-oss-20b":                    "GPT-OSS 20B (OpenAI)",
    # Preview models (may be discontinued at short notice)
    "meta-llama/llama-4-scout-17b-16e-instruct": "Llama 4 Scout 17B (preview)",
    "qwen/qwen3-32b":                        "Qwen3 32B (preview)",
}

DEFAULT_GROQ_MODEL = "llama-3.3-70b-versatile"

def get_llm(model_override: str = None):
    """
    Returns a Groq-backed LLM.
    Base model: llama-3.3-70b-versatile (set via GROQ_MODEL env var).
    model_override accepts any key from GROQ_MODELS or any raw Groq model id.
    DeepSeek models are served by Groq (deepseek-r1-distill-llama-70b etc.)
    so no Ollama installation is required.
    """
    load_dotenv(override=True)
    groq_api_key = os.getenv("GROQ_API_KEY")

    if not groq_api_key or not groq_api_key.strip():
        raise RuntimeError("[LLM] ❌ GROQ_API_KEY is not set in .env")

    groq_model = model_override if model_override else os.getenv("GROQ_MODEL", DEFAULT_GROQ_MODEL)

    # Validate — warn but don't block for unknown model ids (Groq adds new models often)
    if groq_model not in GROQ_MODELS:
        print(f"[LLM] ⚠️  Unknown model id '{groq_model}' — sending to Groq anyway.")
    else:
        print(f"[LLM] Using Groq model: {groq_model} ({GROQ_MODELS[groq_model]})")

    return ChatGroq(
        api_key=groq_api_key,
        model=groq_model,
        temperature=0,
    )

def delete_file_chunks(file_name: str, course_id: str = None):
    """
    Deletes existing chunks for a specific file from the vector store to prevent duplication.
    ChromaDB requires the $and operator when filtering on more than one field.
    """
    try:
        vs = get_vector_store()

        # ChromaDB only allows a single key at the top level of `where`.
        # When we have multiple conditions we must use the $and operator.
        if course_id:
            where = {"$and": [{"source": file_name}, {"course_id": course_id}]}
        else:
            where = {"source": file_name}

        results = vs.get(where=where)
        ids_to_delete = results.get("ids", [])

        if ids_to_delete:
            vs.delete(ids_to_delete)
            print(f"Deleted {len(ids_to_delete)} existing chunks for {file_name}")
    except Exception as e:
        print(f"Error during deduplication: {e}")

def _load_pptx(file_path: str, course_id: str = None, user_id: str = None) -> list:
    """
    Extract text AND images from a .pptx file using python-pptx (no LibreOffice needed).
    Returns a list of LangChain Document objects:
      - One text Document per slide (existing behaviour)
      - One image-caption Document per image/diagram (new — via Groq Vision)
    """
    from pptx import Presentation
    prs = Presentation(file_path)
    source_name = os.path.basename(file_path)
    text_docs = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        content = "\n".join(texts)
        if content.strip():
            text_docs.append(Document(
                page_content=content,
                metadata={"page": slide_num, "source": source_name}
            ))

    print(f"[_load_pptx] Extracted {len(text_docs)} slide(s) of text from '{source_name}'")

    # ── Vision: extract image captions from each slide ───────────────────────
    image_docs = []
    try:
        from vision_extractor import process_pptx_images
        image_docs = process_pptx_images(
            pptx_path=file_path,
            source_name=source_name,
            course_id=course_id,
            user_id=user_id,
        )
        if image_docs:
            print(f"[_load_pptx] 🖼️  Vision added {len(image_docs)} image caption(s) from '{source_name}'")
    except Exception as e:
        print(f"[_load_pptx] ⚠️  Vision extraction failed (skipping): {e}")
    # ─────────────────────────────────────────────────────────────────────────

    return text_docs + image_docs


def _load_ppt_via_com(file_path: str) -> list:
    """
    Convert legacy .ppt (binary OLE format) to .pptx using PowerPoint COM automation
    (Windows only — uses the already-installed Microsoft PowerPoint, no LibreOffice needed),
    then load it with _load_pptx.
    """
    try:
        import comtypes.client  # ships with pywin32 / comtypes, already on Windows
        print(f"[_load_ppt_via_com] Converting '{os.path.basename(file_path)}' via PowerPoint COM...")

        # Create a temp .pptx path
        tmp_pptx = tempfile.mktemp(suffix=".pptx")
        abs_path = os.path.abspath(file_path)

        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1  # Must be visible on some Windows setups
        deck = powerpoint.Presentations.Open(abs_path, ReadOnly=True, Untitled=True, WithWindow=False)
        deck.SaveAs(tmp_pptx, 24)  # 24 = ppSaveAsOpenXMLPresentation (.pptx)
        deck.Close()
        powerpoint.Quit()

        print(f"[_load_ppt_via_com] Converted to temp .pptx: {tmp_pptx}")
        docs = _load_pptx(tmp_pptx)

        # Cleanup temp file
        try:
            os.remove(tmp_pptx)
        except Exception:
            pass

        return docs

    except ImportError:
        raise RuntimeError(
            "Could not import 'comtypes'. Install it with:  pip install comtypes\n"
            "Or convert your .ppt file to .pptx manually and re-upload."
        )
    except Exception as e:
        raise RuntimeError(
            f"Failed to convert '{os.path.basename(file_path)}' using PowerPoint COM.\n"
            f"Tip: Make sure Microsoft PowerPoint is installed, then try again.\n"
            f"Or convert the file to .pptx manually and re-upload.\n"
            f"Original error: {e}"
        )


def ingest_file_from_path(file_path: str, course_id: str = None, user_id: str = None):
    """
    Ingests a file (.pdf, .docx, .pptx, .ppt) from a local file path with deduplication.
    Now also extracts and AI-captions images/diagrams from PDFs and PPTX via the Vision LLM.
    """
    try:
        file_name = os.path.basename(file_path)
        ext = file_path.lower().split('.')[-1]

        if ext == 'pdf':
            loader = PyPDFLoader(file_path)
            docs = loader.load()
            # ── Vision: extract and caption images from each PDF page ────────
            try:
                from vision_extractor import process_pdf_images
                image_docs = process_pdf_images(
                    pdf_path=file_path,
                    source_name=file_name,
                    course_id=course_id,
                    user_id=user_id,
                )
                if image_docs:
                    print(f"[RAG] 🖼️  Vision added {len(image_docs)} image caption(s) from '{file_name}'")
                    docs = docs + image_docs
            except Exception as e:
                print(f"[RAG] ⚠️  PDF vision extraction failed (skipping): {e}")
            # ─────────────────────────────────────────────────────────────────
        elif ext == 'docx':
            loader = Docx2txtLoader(file_path)
            docs = loader.load()
            # Fallback: Docx2txtLoader sometimes returns empty content for
            # certain .docx files (e.g. text inside shapes/tables only).
            # Try python-docx directly in that case.
            if not docs or not any(d.page_content.strip() for d in docs):
                print(f"[RAG] Docx2txtLoader returned no text for '{file_name}', trying python-docx fallback...")
                try:
                    from docx import Document as DocxDocument
                    docx_doc = DocxDocument(file_path)
                    full_text = "\n".join(
                        para.text for para in docx_doc.paragraphs if para.text.strip()
                    )
                    if full_text.strip():
                        from langchain_core.documents import Document
                        docs = [Document(page_content=full_text, metadata={"source": file_name, "page": 1})]
                    else:
                        print(f"[RAG] python-docx also found zero text for '{file_name}'.")
                except ImportError:
                    print(f"[RAG] python-docx not installed, fallback skipped for '{file_name}'.")
        elif ext == 'pptx':
            # Pass course_id/user_id so vision_extractor can tag image docs correctly
            docs = _load_pptx(file_path, course_id=course_id, user_id=user_id)
        elif ext == 'ppt':
            docs = _load_ppt_via_com(file_path)
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ".", " ", ""]
        )
        splits = text_splitter.split_documents(docs)

        # Handle completely empty files (e.g. image-only docx files) instead of crashing
        if not splits:
            print(f"[RAG] Warning: '{file_name}' produced 0 text chunks. Inserting placeholder.")
            from langchain_core.documents import Document
            placeholder_text = f"The document '{file_name}' was uploaded but contained no extractable text. It may be an image-only scan or format not supported by the text extractor."
            
            # Create a single placeholder chunk
            dummy_doc = Document(
                page_content=placeholder_text,
                metadata={
                    "source": file_name,
                    "page": 1,
                    "course_id": course_id or "unknown",
                    "user_id": user_id or "unknown"
                }
            )
            splits = [dummy_doc]

        # Add metadata
        for split in splits:
            if course_id:
                split.metadata["course_id"] = course_id
            if user_id:
                split.metadata["user_id"] = user_id
            split.metadata["source"] = file_name

        # Deduplicate before adding
        delete_file_chunks(file_name, course_id)

        # Add to vector store
        vs = get_vector_store()
        vs.add_documents(splits)
        print(f"[RAG] ✅ Ingested '{file_name}': {len(splits)} chunks added.")
        return {"status": "success", "chunks_added": len(splits)}
    except Exception as e:
        print(f"Error ingesting {file_path}: {e}")
        raise e

def ingest_file_from_bytes(file_bytes: bytes, file_name: str, course_id: str = None, user_id: str = None):
    """
    Ingests a file from raw bytes (e.g. fetched in-memory from Google Drive).
    Creates a temp directory, saves the file with its REAL name, ingests, then cleans up.
    This ensures the 'source' metadata in ChromaDB matches the original filename.
    No permanent disk storage needed.
    """
    import tempfile
    import shutil
    tmp_dir = None
    try:
        tmp_dir = tempfile.mkdtemp()
        tmp_path = os.path.join(tmp_dir, file_name)
        with open(tmp_path, 'wb') as f:
            f.write(file_bytes)
        print(f"[RAG] Ingesting '{file_name}' from memory via temp file...")
        result = ingest_file_from_path(tmp_path, course_id=course_id, user_id=user_id)
        return result
    except Exception as e:
        print(f"[RAG] Error ingesting '{file_name}' from bytes: {e}")
        raise e
    finally:
        if tmp_dir and os.path.exists(tmp_dir):
            try:
                shutil.rmtree(tmp_dir)
            except Exception:
                print(f"[RAG] Warning: Could not clean up temp dir {tmp_dir}")



# Deprecated/Alias for backward compatibility
def ingest_pdf_from_path(file_path: str, course_id: str = None, user_id: str = None):
    return ingest_file_from_path(file_path, course_id, user_id)

def ingest_document(file_upload, user_id: str = None):
    """
    Ingests a PDF document from an uploaded file.
    """
    # Determine extension
    ext = ".pdf"
    if hasattr(file_upload, 'filename'):
        filename = file_upload.filename
        if '.' in filename:
            ext = '.' + filename.split('.')[-1]

    # Save uploaded file temporarily
    with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        try:
            tmp.write(file_upload.file.read())
            tmp_path = tmp.name
            return ingest_file_from_path(tmp_path, user_id=user_id)
        finally:
            try:
                # Ensure the file is closed before trying to remove it
                pass
                os.unlink(tmp_path)
            except PermissionError:
                # On Windows, sometimes file is still locked.
                print(f"Warning: Could not delete temp file {tmp_path}")

def query_rag_stream(question: str, marks: int = 5, course_id: str = None, user_id: str = None, model: str = None):
    """
    Queries the RAG system and yields tokens as they are generated by the LLM.
    """
    SEP = "=" * 70
    print(f"\n{SEP}")
    print(f"[RAG] 🔍 QUERY RECEIVED")
    print(f"  Question : {question}")
    print(f"  Marks    : {marks}")
    print(f"  Course ID: {course_id}")
    print(f"  Model    : {model or f'default ({DEFAULT_GROQ_MODEL})'}")
    print(SEP)

    print(f"[RAG] ⚙️  Initializing LLM ({model or 'default'})...")
    llm = get_llm(model_override=model)
    print(f"[RAG] ✅ LLM ready. Now retrieving documents from ChromaDB...")
    
    # Configure MMR (Maximal Marginal Relevance) for better diversity and relevance
    search_kwargs = {"k": 10, "fetch_k": 30}
    
    # Construct filter for ChromaDB
    filters = {}
    if course_id:
        filters["course_id"] = course_id
    if user_id:
        filters["user_id"] = user_id
        
    if filters:
        if len(filters) > 1:
            # Chroma handles $and naturally or can just pass the dict directly
            search_kwargs["filter"] = {"$and": [{k: v} for k, v in filters.items()]}
        else:
            search_kwargs["filter"] = filters
        
    retriever = get_vector_store().as_retriever(
        search_type="mmr",
        search_kwargs=search_kwargs
    )

    # ── DEBUG: show retrieved chunks before building the prompt ──────────────
    _raw_docs = retriever.invoke(question)
    print(f"\n[RAG] 📄 RETRIEVED {len(_raw_docs)} CHUNK(S) from ChromaDB")
    print("-" * 70)
    for i, doc in enumerate(_raw_docs, start=1):
        src  = doc.metadata.get("source", "unknown")
        page = doc.metadata.get("page", "?")
        cid  = doc.metadata.get("course_id", "—")
        snippet = doc.page_content.replace("\n", " ").strip()[:200]
        print(f"  Chunk {i}/{len(_raw_docs)}")
        print(f"    Source    : {src}")
        print(f"    Page      : {page}")
        print(f"    Course ID : {cid}")
        print(f"    Snippet   : {snippet}...")
        print()
    if not _raw_docs:
        print("  ⚠️  No chunks found — LLM will respond with 'cannot answer'.")
    print("-" * 70)
    # ────────────────────────────────────────────────────────────────────────

    # ── Build marks-specific instructions ───────────────────────────────────
    format_instruction = ""
    marks_val = str(marks).strip()
    marks_int = int(marks_val) if marks_val in ["1", "2", "5", "10"] else 0

    if marks_val in ["2", "5", "10", 2, 5, 10]:
        marks_int = int(marks_val)
        if marks_int == 2:
            length_instruction = "Answer in EXACTLY 40 to 50 words. Plain text only — no code blocks, no diagrams."
            format_instruction = """
Write a clean 40-50 word text answer. Use this format:
**Definition:** One clear definition sentence.
**Key Points:** 2-3 brief bullet points.
**Application:** Brief example if needed.

STRICT RULE for 2 marks: Do NOT include any code blocks, image descriptions, or lengthy explanations."""
        elif marks_int == 5:
            length_instruction = "Answer in 20 to 25 lines. Cover the topic comprehensively."
            format_instruction = """
Use this exact exam format. Each section MUST start on a NEW LINE:
**Topic:** name of topic
**Definition:** 3-4 line definition
**Types/Categories:**
- Type 1: description
- Type 2: description
**Working/Explanation:** 5-6 lines explaining how it works
**Conclusion:** 1-2 lines

VISUAL & CODE SUPPLEMENT (for 5-mark answers):
- If the Context contains any actual code, reproduce it EXACTLY under a heading **Code Example:** in a fenced code block. If the Context discusses a coding topic without a snippet, write a short illustrative code block yourself.
- If the Context contains an IMAGE DESCRIPTION (diagrams, flowcharts, graphs) relevant to this topic, you MUST extract and describe it clearly under a heading **Diagram/Visual:**.
- These sections are ADDITIONAL to the main text.

IMPORTANT: Every heading and bullet MUST be on its own line. Do NOT merge sections."""
        elif marks_int == 10:
            length_instruction = "Answer in 40 to 45 lines. Provide an extremely detailed and comprehensive explanation."
            format_instruction = """
Use this exact exam format. Every section and bullet MUST be on its own NEW LINE:
**Topic:** name of topic
**Introduction & Definition:** 4-5 line detailed definition and background
**Types/Classification:**
1. Type 1: detailed description
2. Type 2: detailed description
**Detailed Explanation/Working:** 10-12 lines of in-depth explanation
**Advantages & Limitations:**
- Advantage 1 (detailed)
- Limitation 1 (detailed)
**Conclusion:** 2-3 concluding lines

VISUAL & CODE SUPPLEMENT (mandatory for 10-mark answers if relevant):
- If the Context contains actual code, pseudocode, or algorithms → reproduce it EXACTLY in a fenced code block under **Code Example:**. If the Context discusses a coding topic without a snippet, write a comprehensive, well-commented code block yourself.
- If the Context contains an IMAGE DESCRIPTION (diagrams, flowcharts, architecture) relevant to the topic → you MUST extract and describe its flow, relationships, and components in detail under a **Diagram/Visual:** heading.
- These sections are MANDATORY for 10-mark answers if the topic allows!

IMPORTANT: Every heading, numbered point, and bullet MUST be on its own line. Do NOT merge sections."""
    elif marks_val == "1" or marks_val == 1:
        length_instruction = "Answer in EXACTLY 2 lines. One definition line and one example/use-case line."
    else:
        length_instruction = "Answer in a moderate paragraph."

    # Warn early if no context was retrieved
    if not _raw_docs:
        print("[RAG] ⚠️  No context retrieved — LLM will be told to refuse the answer.")

    template = f"""You are an expert academic exam assistant helping students study. Answer questions using the provided study materials below.

CRITICAL RULES:
1. Base your factual text answer on the content from the Context section below.
2. Do NOT fabricate historical facts, names, dates, or textual definitions that are NOT in the Context.
3. If the Context does not mention the topic at all, respond ONLY with: "I cannot answer this based on the provided materials."
4. At the end of every answer you MUST cite source file(s). Example: "Source: lecture1.pdf, Page 5"
5. If writing code examples for 5/10 marks, you may use external knowledge to write syntactically correct code that illustrates the concepts discussed in the Context.

Context (primary source of truth):
{{context}}

Question: {{question}}

Marks: {marks_int} marks
Instruction: {length_instruction}
{format_instruction}"""
    
    prompt = ChatPromptTemplate.from_template(template)
    
    def format_docs(docs):
        """Format retrieved docs into a clearly labelled context string for the LLM."""
        if not docs:
            return "[No relevant context found in the uploaded materials.]"
        lines = []
        for doc in docs:
            source = doc.metadata.get("source", "Unknown Document")
            page   = doc.metadata.get("page",   "?")
            lines.append(
                f"--- BEGIN EXCERPT [File: {source}, Page: {page}] ---\n"
                f"{doc.page_content.strip()}\n"
                f"--- END EXCERPT ---\n"
            )
        return "\n".join(lines)
    
    chain = (
        {"context": retriever | format_docs, "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

    # ── DEBUG: stream tokens and print each one to terminal ─────────────────
    print(f"[RAG] 🤖 LLM STREAMING — tokens below:")
    print("-" * 70)
    token_count = 0

    def _stream_chain(active_chain):
        """Inner generator to stream from the given chain."""
        for chunk in active_chain.stream(question):
            print(chunk, end="", flush=True)
            yield chunk

        filename = file_upload.filename
        if '.' in filename:
            ext = '.' + filename.split('.')[-1]

    # Save uploaded file temporarily
    with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        try:
            tmp.write(file_upload.file.read())
            tmp_path = tmp.name
            return ingest_file_from_path(tmp_path, user_id=user_id)
        finally:
            try:
                # Ensure the file is closed before trying to remove it
                pass
                os.unlink(tmp_path)
            except PermissionError:
                # On Windows, sometimes file is still locked.
                print(f"Warning: Could not delete temp file {tmp_path}")

def query_rag_stream(question: str, marks: int = 5, course_id: str = None, user_id: str = None, model: str = None):
    """
    Queries the RAG system and yields tokens as they are generated by the LLM.
    """
    SEP = "=" * 70
    print(f"\n{SEP}")
    print(f"[RAG] 🔍 QUERY RECEIVED")
    print(f"  Question : {question}")
    print(f"  Marks    : {marks}")
    print(f"  Course ID: {course_id}")
    print(f"  Model    : {model or f'default ({DEFAULT_GROQ_MODEL})'}")
    print(SEP)

    print(f"[RAG] ⚙️  Initializing LLM ({model or 'default'})...")
    llm = get_llm(model_override=model)
    print(f"[RAG] ✅ LLM ready. Now retrieving documents from ChromaDB...")
    
    # Configure MMR (Maximal Marginal Relevance) for better diversity and relevance
    search_kwargs = {"k": 10, "fetch_k": 30}
    
    # Construct filter for ChromaDB
    filters = {}
    if course_id:
        filters["course_id"] = course_id
    if user_id:
        filters["user_id"] = user_id
        
    if filters:
        if len(filters) > 1:
            # Chroma handles $and naturally or can just pass the dict directly
            search_kwargs["filter"] = {"$and": [{k: v} for k, v in filters.items()]}
        else:
            search_kwargs["filter"] = filters
        
    retriever = get_vector_store().as_retriever(
        search_type="mmr",
        search_kwargs=search_kwargs
    )

    # ── DEBUG: show retrieved chunks before building the prompt ──────────────
    _raw_docs = retriever.invoke(question)
    print(f"\n[RAG] 📄 RETRIEVED {len(_raw_docs)} CHUNK(S) from ChromaDB")
    print("-" * 70)
    for i, doc in enumerate(_raw_docs, start=1):
        src  = doc.metadata.get("source", "unknown")
        page = doc.metadata.get("page", "?")
        cid  = doc.metadata.get("course_id", "—")
        snippet = doc.page_content.replace("\n", " ").strip()[:200]
        print(f"  Chunk {i}/{len(_raw_docs)}")
        print(f"    Source    : {src}")
        print(f"    Page      : {page}")
        print(f"    Course ID : {cid}")
        print(f"    Snippet   : {snippet}...")
        print()
    if not _raw_docs:
        print("  ⚠️  No chunks found — LLM will respond with 'cannot answer'.")
    print("-" * 70)
    # ────────────────────────────────────────────────────────────────────────

    # ── Build marks-specific instructions ───────────────────────────────────
    format_instruction = ""
    marks_val = str(marks).strip()
    marks_int = int(marks_val) if marks_val in ["1", "2", "5", "10"] else 0

    if marks_val in ["2", "5", "10", 2, 5, 10]:
        marks_int = int(marks_val)
        if marks_int == 2:
            length_instruction = "Answer in EXACTLY 40 to 50 words. Plain text only — no code blocks, no diagrams."
            format_instruction = """
Write a clean 40-50 word text answer. Use this format:
**Definition:** One clear definition sentence.
**Key Points:** 2-3 brief bullet points.
**Application:** Brief example if needed.

STRICT RULE for 2 marks: Do NOT include any code blocks, image descriptions, or lengthy explanations."""
        elif marks_int == 5:
            length_instruction = "Answer in 20 to 25 lines. Cover the topic comprehensively."
            format_instruction = """
Use this exact exam format. Each section MUST start on a NEW LINE:
**Topic:** name of topic
**Definition:** 3-4 line definition
**Types/Categories:**
- Type 1: description
- Type 2: description
**Working/Explanation:** 5-6 lines explaining how it works
**Conclusion:** 1-2 lines

VISUAL & CODE SUPPLEMENT (for 5-mark answers):
- If the Context contains any actual code, reproduce it EXACTLY under a heading **Code Example:** in a fenced code block. If the Context discusses a coding topic without a snippet, write a short illustrative code block yourself.
- If the Context contains an IMAGE DESCRIPTION (diagrams, flowcharts, graphs) relevant to this topic, you MUST extract and describe it clearly under a heading **Diagram/Visual:**.
- These sections are ADDITIONAL to the main text.

IMPORTANT: Every heading and bullet MUST be on its own line. Do NOT merge sections."""
        elif marks_int == 10:
            length_instruction = "Answer in 40 to 45 lines. Provide an extremely detailed and comprehensive explanation."
            format_instruction = """
Use this exact exam format. Every section and bullet MUST be on its own NEW LINE:
**Topic:** name of topic
**Introduction & Definition:** 4-5 line detailed definition and background
**Types/Classification:**
1. Type 1: detailed description
2. Type 2: detailed description
**Detailed Explanation/Working:** 10-12 lines of in-depth explanation
**Advantages & Limitations:**
- Advantage 1 (detailed)
- Limitation 1 (detailed)
**Conclusion:** 2-3 concluding lines

VISUAL & CODE SUPPLEMENT (mandatory for 10-mark answers if relevant):
- If the Context contains actual code, pseudocode, or algorithms → reproduce it EXACTLY in a fenced code block under **Code Example:**. If the Context discusses a coding topic without a snippet, write a comprehensive, well-commented code block yourself.
- If the Context contains an IMAGE DESCRIPTION (diagrams, flowcharts, architecture) relevant to the topic → you MUST extract and describe its flow, relationships, and components in detail under a **Diagram/Visual:** heading.
- These sections are MANDATORY for 10-mark answers if the topic allows!

IMPORTANT: Every heading, numbered point, and bullet MUST be on its own line. Do NOT merge sections."""
    elif marks_val == "1" or marks_val == 1:
        length_instruction = "Answer in EXACTLY 2 lines. One definition line and one example/use-case line."
    else:
        length_instruction = "Answer in a moderate paragraph."

    # Warn early if no context was retrieved
    if not _raw_docs:
        print("[RAG] ⚠️  No context retrieved — LLM will be told to refuse the answer.")

    template = f"""You are an expert academic exam assistant helping students study. Answer questions using the provided study materials below.

CRITICAL RULES:
1. Base your factual text answer on the content from the Context section below.
2. Do NOT fabricate historical facts, names, dates, or textual definitions that are NOT in the Context.
3. If the Context does not mention the topic at all, respond ONLY with: "I cannot answer this based on the provided materials."
4. At the end of every answer you MUST cite source file(s). Example: "Source: lecture1.pdf, Page 5"
5. If writing code examples for 5/10 marks, you may use external knowledge to write syntactically correct code that illustrates the concepts discussed in the Context.

Context (primary source of truth):
{{context}}

Question: {{question}}

Marks: {marks_int} marks
Instruction: {length_instruction}
{format_instruction}"""
    
    prompt = ChatPromptTemplate.from_template(template)
    
    def format_docs(docs):
        """Format retrieved docs into a clearly labelled context string for the LLM."""
        if not docs:
            return "[No relevant context found in the uploaded materials.]"
        lines = []
        for doc in docs:
            source = doc.metadata.get("source", "Unknown Document")
            page   = doc.metadata.get("page",   "?")
            lines.append(
                f"--- BEGIN EXCERPT [File: {source}, Page: {page}] ---\n"
                f"{doc.page_content.strip()}\n"
                f"--- END EXCERPT ---\n"
            )
        return "\n".join(lines)
    
    chain = (
        {"context": retriever | format_docs, "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

    # ── DEBUG: stream tokens and print each one to terminal ─────────────────
    print(f"[RAG] 🤖 LLM STREAMING — tokens below:")
    print("-" * 70)
    token_count = 0

    def _stream_chain(active_chain):
        """Inner generator to stream from the given chain."""
        for chunk in active_chain.stream(question):
            print(chunk, end="", flush=True)
            yield chunk

    try:
        for chunk in _stream_chain(chain):
            token_count += 1
            yield chunk
    except Exception as stream_err:
        print(f"\n[RAG] ❌ GROQ ERROR: {stream_err}")
        print(f"[RAG] ❌ Error type: {type(stream_err).__name__}")
        raise

    print(f"\n{'-' * 70}")
    print(f"[RAG] ✅ Stream complete. Total token chunks emitted: {token_count}")
    print("=" * 70 + "\n")
    # ────────────────────────────────────────────────────────────────────────

def query_rag(question: str, marks: int = 5, course_id: str = None, user_id: str = None):
    """
    Non-streaming fallback (deprecated).
    """
    response = ""
    for chunk in query_rag_stream(question, marks, course_id, user_id):
        response += chunk
    return response

def delete_course_data(course_id: str):
    """
    Deletes all data associated with a specific course_id.
    """
    try:
        print(f"Deleting data for course: {course_id}")
        vs = get_vector_store()
        
        results = vs.get(where={"course_id": course_id})
        ids_to_delete = results.get("ids", [])
        
        if ids_to_delete:
            vs.delete(ids_to_delete)
            print(f"Deleted {len(ids_to_delete)} chunks for course {course_id}")
        else:
            print(f"No chunks found for course {course_id}")
            
        return True
    except Exception as e:
        print(f"Error deleting course data: {e}")
        return False

def generate_assessment(topic: str, course_id: str, assessment_type: str, count: int = 5, model: str = None, file_name: str = None) -> str:
    """
    Generates an assessment (quiz or flashcards) in strict JSON format based on course material.
    """
    llm = get_llm(model)
    
    # Force JSON mode for models that support it
    # Note: Groq supports response_format={"type": "json_object"}
    try:
        llm = llm.bind(response_format={"type": "json_object"})
    except Exception:
        pass # fallback if model binding fails

    # Retrieve context with MMR (Maximal Marginal Relevance) 
    # to ensure diversity, pulling chunks from across the entire document
    search_kwargs = {"k": 20, "fetch_k": 100}
    
    # Build filter conditions
    if course_id and file_name:
        if file_name == "ALL":
            search_kwargs["filter"] = {"course_id": course_id}
        else:
            search_kwargs["filter"] = {"$and": [{"course_id": course_id}, {"source": file_name}]}
    elif course_id:
        search_kwargs["filter"] = {"course_id": course_id}
    elif file_name and file_name != "ALL":
        search_kwargs["filter"] = {"source": file_name}
        
    retriever = get_vector_store().as_retriever(
        search_type="mmr",
        search_kwargs=search_kwargs
    )
    
    docs = retriever.invoke(topic)
    
    if not docs:
        if assessment_type == "quiz":
            return '{"questions": []}'
        else:
            return '{"flashcards": []}'
            
    context_text = "\n".join([d.page_content for d in docs])

    if assessment_type == "quiz":
        instruction = f"""Generate a {count}-question multiple-choice quiz about "{topic}" based strictly on the provided Context.
You MUST output valid, parseable JSON exactly matching this schema:
{{
  "questions": [
    {{
      "question": "The question text?",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correctIndex": 0,
      "explanation": "A short explanation of why this is correct based on the context."
    }}
  ]
}}
Ensure "correctIndex" is an integer between 0 and 3."""
    elif assessment_type == "flashcard":
        instruction = f"""Generate {count} flashcards about "{topic}" based strictly on the provided Context.
You MUST output valid, parseable JSON exactly matching this schema:
{{
  "flashcards": [
    {{
      "front": "A clear question or concept name?",
      "back": "A concise, accurate definition or explanation."
    }}
  ]
}}"""
    else:
        return "{}"

    template = f"""You are an educational AI. Your ONLY output must be a valid JSON object.
Do NOT wrap the JSON in markdown blocks like ```json ... ```. Output raw JSON only.

Context from Course Materials:
{context_text}

{instruction}
"""
    try:
        from langchain_core.messages import SystemMessage, HumanMessage
        messages = [
            SystemMessage(content="You are a JSON-only API. You must return valid JSON without any markdown formatting, preamble, or postscript."),
            HumanMessage(content=template)
        ]
        response = llm.invoke(messages)
        content = response.content.strip()
        
        # Cleanup in case the LLM ignored instructions and returned markdown blocks
        if content.startswith("```json"):
            content = content[7:]
        if content.startswith("```"):
            content = content[3:]
        if content.endswith("```"):
            content = content[:-3]
            
        return content.strip()
    except Exception as e:
        print(f"Error generating assessment: {e}")
        if assessment_type == "quiz":
            return '{"questions": []}'
        else:
            return '{"flashcards": []}'
